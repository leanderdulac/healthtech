#!/usr/bin/env python3
"""Gera apresentação PPTX com os diagramas do projeto Healthtech."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).parent
DIAGRAMS = DOCS / "diagrams"
OUTPUT = DOCS / "Healthtech_Datalake_VertexAI.pptx"

NAVY = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x02, 0x80, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x29, 0x5C)
LIGHT = RGBColor(0xF2, 0xF7, 0xFA)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(8.8), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)
    p2.space_before = Pt(12)


def add_diagram_slide(prs: Presentation, title: str, image_path: Path, notes: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
    header.fill.solid()
    header.fill.fore_color.rgb = TEAL
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(9), Inches(0.5))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    slide.shapes.add_picture(str(image_path), Inches(0.3), Inches(0.85), width=Inches(9.4))

    if notes:
        note_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.15), Inches(9.4), Inches(0.4))
        np_ = note_box.text_frame.paragraphs[0]
        np_.text = notes
        np_.font.size = Pt(11)
        np_.font.color.rgb = DARK


def add_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "Entry Points e Artefatos"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = NAVY

    items = [
        "run_datalake_pipeline.py — Datalake Bronze → Silver → Gold",
        "run_vertex_integration.py — Datalake + Treino + Online + Batch",
        "main_simulation.py — Demo original com Vertex isolado",
        "data/lakehouse/ — Parquet particionado (Bronze/Silver/Gold)",
        "data/vertex_exports/ — JSONL e CSV para Vertex AI",
        "data/models/ — Isolation Forest local (fallback)",
        "docs/diagrams/ — Fontes Mermaid + PNG + SVG",
    ]

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(4))
    tf = body.text_frame
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        p.level = 0


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(
        prs,
        "Healthtech — Datalake de Telemetria 24h",
        "Arquitetura Medallion + Integração Vertex AI\nWearables · FHIR · Early Warning",
    )

    slides = [
        ("1. Visão Geral — End-to-End", "01-visao-geral.png",
         "Fluxo completo: wearables → datalake → extração → Vertex AI"),
        ("2. Sequência de Execução", "02-sequencia.png",
         "Ordem temporal do run_vertex_integration.py"),
        ("3. Camadas Medallion", "03-medallion.png",
         "Bronze (raw) → Silver (curated) → Gold (analytics)"),
        ("4. Modos Vertex AI", "04-vertex-modos.png",
         "Produção GCP vs modo local (Isolation Forest + heurísticas)"),
    ]

    for title, img, note in slides:
        add_diagram_slide(prs, title, DIAGRAMS / img, note)

    add_summary_slide(prs)

    closing = prs.slides.add_slide(prs.slide_layouts[6])
    closing.background.fill.solid()
    closing.background.fill.fore_color.rgb = NAVY
    cb = closing.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    cp = cb.text_frame.paragraphs[0]
    cp.text = "Próximo passo: configurar .env + gcloud auth\npara ativar Vertex AI em produção"
    cp.font.size = Pt(24)
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER

    prs.save(str(OUTPUT))
    print(f"Apresentação gerada: {OUTPUT}")


if __name__ == "__main__":
    main()